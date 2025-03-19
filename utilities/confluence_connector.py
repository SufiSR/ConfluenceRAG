import logging
from typing import Dict, List, Tuple, Set
import tempfile
import os
import docx2txt
import fitz  # PyMuPDF for PDF processing
import pptx
from bs4 import BeautifulSoup
import time

from configuration import get_confluence_client

logger = logging.getLogger(__name__)


class ConfluenceConnector:
    def __init__(self, profile="internal-confluence"):
        self.client = get_confluence_client(profile)
        self.profile = profile

    def get_all_spaces(self) -> List[Dict]:
        """Get all available spaces."""
        return self.client.get_all_spaces()

    def get_pages_from_space(self, space_key: str) -> List[Dict]:
        """Get all pages from a specific space."""
        pages = []
        start = 0
        limit = 100

        while True:
            results = self.client.get_all_pages_from_space(space_key, start=start, limit=limit)
            if not results:
                break
            pages.extend(results)
            start += limit
            if len(results) < limit:
                break

        return pages

    def get_page_content(self, page_id: str) -> Tuple[str, Dict]:
        """Get the content and metadata of a specific page."""
        page = self.client.get_page_by_id(page_id, expand='body.storage,version')

        if not page:
            return "", {}

        content = page.get('body', {}).get('storage', {}).get('value', '')

        # Parse HTML content to extract text
        soup = BeautifulSoup(content, 'html.parser')
        text = soup.get_text(strip=True)

        # Extract space key
        space_key = page.get('space', {}).get('key', '')

        # Extract metadata
        metadata = {
            'title': page.get('title', ''),
            # Fix link format
            'link': f"{self.client.url}/wiki/spaces/{space_key}/pages/{page_id}",
            'last_modified': page.get('version', {}).get('when', ''),
            'source_url': f"{self.client.url}/wiki/rest/api/content/{page_id}",
            'space': space_key,
            'page_id': page_id,
        }

        return text, metadata

    def get_page_hierarchy(self, page_id: str) -> List[str]:
        """Get the hierarchy of a page (breadcrumbs)."""
        ancestors = self.client.get_page_ancestors(page_id)
        return [ancestor.get('title', '') for ancestor in ancestors]

    def get_child_pages(self, page_id: str) -> List[Dict]:
        """Get all child pages of a specific page."""
        return self.client.get_page_child_by_type(page_id, type='page')

    def get_attachments(self, page_id: str) -> List[Dict]:
        """Get all attachments of a specific page."""
        return self.client.get_attachments_from_content(page_id)

    def get_attachment_content(self, attachment: Dict, parent_page_id: str) -> Tuple[str, Dict]:
        """Extract text from an attachment based on its file type."""
        attachment_id = attachment.get('id', '')
        filename = attachment.get('title', '')
        file_extension = os.path.splitext(filename)[1].lower()

        # Skip unsupported file types
        if file_extension not in ['.pdf', '.docx', '.pptx']:
            return "", {}

        try:
            # Get attachment download URL
            download_url = attachment.get('_links', {}).get('download', '')

            if not download_url:
                logger.error(f"No download URL found for attachment {filename}")
                return "", {}

            # Download attachment content
            # Use the full URL if it's absolute, otherwise append to base URL
            if download_url.startswith('http'):
                full_url = download_url
            else:
                full_url = f"{self.client.url}{download_url}"

            # Use the requests session from the client to download
            response = self.client._session.get(full_url)

            if response.status_code != 200:
                logger.error(f"Failed to download attachment {filename}: HTTP {response.status_code}")
                return "", {}

            attachment_content = response.content

            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
                temp_file.write(attachment_content)
                temp_file_path = temp_file.name

            text = ""

            # Extract text based on file type
            if file_extension == '.pdf':
                text = self._extract_text_from_pdf(temp_file_path)
            elif file_extension == '.docx':
                text = self._extract_text_from_docx(temp_file_path)
            elif file_extension == '.pptx':
                text = self._extract_text_from_pptx(temp_file_path)

            # Clean up temporary file
            os.unlink(temp_file_path)

            # Create metadata
            metadata = {
                'title': filename,
                'link': full_url,
                'last_modified': attachment.get('version', {}).get('when', ''),
                'source_url': attachment.get('_links', {}).get('self', ''),
                'space': attachment.get('space', ''),
                'parent_page_id': parent_page_id,
                'attachment_id': attachment_id,
                'is_attachment': True,
                'file_type': file_extension[1:],  # Remove the dot
            }

            return text, metadata

        except Exception as e:
            logger.error(f"Error extracting content from attachment {filename}: {str(e)}")
            return "", {}

    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file."""
        text = ""
        try:
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
        return text

    def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from a DOCX file."""
        try:
            return docx2txt.process(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return ""

    def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from a PPTX file."""
        text = ""
        try:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
        return text

    def get_all_content_from_spaces(self, space_keys: List[str]) -> List[Tuple[str, Dict]]:
        """Get all content (pages and attachments) from specified spaces."""
        all_content = []

        for space_key in space_keys:
            logger.info(f"Processing space: {space_key}")

            # Get all pages from the space
            pages = self.get_pages_from_space(space_key)
            logger.info(f"Found {len(pages)} pages in space {space_key}")

            for page in pages:
                page_id = page.get('id', '')

                # Get page content and metadata
                text, metadata = self.get_page_content(page_id)

                # Add hierarchy information
                hierarchy = self.get_page_hierarchy(page_id)
                metadata['hierarchy'] = hierarchy

                all_content.append((text, metadata))

                # Get attachments
                attachments = self.get_attachments(page_id)

                # Handle the case where attachments is a dictionary with 'results' key
                attachment_list = []
                if isinstance(attachments, dict) and 'results' in attachments:
                    attachment_list = attachments['results']
                elif isinstance(attachments, list):
                    attachment_list = attachments

                logger.info(f"Found {len(attachment_list)} attachments for page {page.get('title', '')}")

                for attachment in attachment_list:
                    # Check if this is a supported file type before processing
                    filename = attachment.get('title', '')
                    file_extension = os.path.splitext(filename)[1].lower()

                    if file_extension not in ['.pdf', '.docx', '.pptx']:
                        continue

                    # Process the attachment
                    attachment_text, attachment_metadata = self.get_attachment_content(attachment, page_id)
                    if attachment_text:
                        # Add hierarchy information to attachment
                        attachment_metadata['hierarchy'] = hierarchy + [page.get('title', '')]
                        all_content.append((attachment_text, attachment_metadata))
                    else:
                        logger.warning(f"Failed to extract text from attachment: {filename}")

        return all_content

    # === New functions for more efficient updates ===
    def get_pages_modified_since(self, space_key: str, modified_since: str) -> List[Dict]:
        """
        Get pages from a specific space that have been modified on or after a given timestamp.
        This uses Confluence's CQL to fetch only the pages that have changed recently.
        Implements a lower pagination limit and retry logic to handle timeouts.

        :param space_key: The key of the space.
        :param modified_since: An ISO-formatted date string.
        :return: A list of pages modified since the provided date.
        """
        global results
        cql = f"space = '{space_key}' AND type = page AND lastmodified >= '{modified_since}'"
        start = 0
        limit = 50  # Reduced limit to reduce load on each request
        pages = []
        max_retries = 3

        while True:
            retries = 0
            while retries < max_retries:
                try:
                    results = self.client.search_content(cql, start=start, limit=limit)
                    break  # If successful, break out of the retry loop
                except Exception as e:
                    retries += 1
                    logger.error(
                        f"Pagination error on space '{space_key}' starting at {start}: {str(e)}. "
                        f"Retry {retries} of {max_retries}."
                    )
                    time.sleep(2 ** retries)  # Exponential backoff

            # If max retries exceeded, break out of the pagination loop
            if retries == max_retries:
                logger.error(
                    f"Exceeded maximum retries for pagination on space '{space_key}' starting at {start}."
                )
                break

            if not results or 'results' not in results:
                break

            pages.extend(results['results'])
            # If the number of results is less than the limit, we've reached the end.
            if len(results['results']) < limit:
                break
            start += limit

        return pages

    def get_page_ids_from_space_minimal(self, space_key: str) -> Set[str]:
        """
        Retrieve a set of page IDs from a specific space using a minimal payload.
        This is useful for quickly identifying which pages still exist in Confluence.
        :param space_key: The key of the space.
        :return: A set of page IDs.
        """
        page_ids = set()
        start = 0
        limit = 100
        try:
            # Try to use a minimal retrieval function if available.
            results = self.client.get_pages_from_space_minimal(space_key, start=start, limit=limit)
        except AttributeError:
            # Fallback: use the existing full retrieval and extract the IDs.
            pages = self.get_pages_from_space(space_key)
            return {page.get('id', '') for page in pages}

        while True:
            if not results:
                break
            for page in results:
                page_ids.add(page.get('id', ''))
            start += limit
            if len(results) < limit:
                break
            results = self.client.get_pages_from_space_minimal(space_key, start=start, limit=limit)
        return page_ids
